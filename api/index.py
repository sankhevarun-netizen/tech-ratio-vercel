"""
Tech Rationalization AI Agent -- Vercel-ready, stateless FastAPI backend.
All state is managed by the client (browser). No server-side sessions.

Local dev:  uvicorn api.index:app --reload --port 8000
Vercel:     vercel deploy
"""

import json, uuid, io, os, re
from typing import Optional, List, Dict, Any
from collections import defaultdict

from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import anthropic
from dotenv import load_dotenv

load_dotenv()

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
ANTHROPIC_MODEL   = os.environ.get("ANTHROPIC_MODEL", "claude-sonnet-4-6")
client = anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY)

# ── FastAPI App ────────────────────────────────────────────────────────────
app = FastAPI(title="Tool Rationalization Agent", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])


# ═══════════════════════════════════════════════════════════════════════════
# SYSTEM PROMPT
# ═══════════════════════════════════════════════════════════════════════════
SYSTEM_PROMPT = """You are an Enterprise Technology Strategy AI Agent specializing in Platform, Application, and Tools Assessment & Rationalization, aligned with KPMG-style consulting frameworks and telecom/large enterprise transformation programs (e.g., Airtel, Tanla-like environments).

You combine the expertise of:
- Senior Enterprise Architect
- Technology Due-Diligence Consultant
- CIO/CTO Advisory Analyst
- Telecom Digital Transformation Expert

PRIMARY OBJECTIVE: Maximize business value, reduce cost, mitigate risk, and simplify the technology landscape by assessing tools across 7 dimensions and applying the 6R Rationalization model.

6R MODEL:
- Retain    -> Strategic, healthy, high-value (Score >=7.5, Risk <=4)
- Rehost    -> Lift-and-shift to cloud (Score >=6, on-prem, cloud-ready)
- Replatform -> Minor modernization (Score 5-7)
- Refactor  -> Significant redesign needed (Score 3-5)
- Replace   -> Better alternative exists (Score <5 or cost-inefficient)
- Retire    -> Decommission -- redundant, EOL, or very low value

SCORING DIMENSIONS (0-10 each):
1. Business Value -- Strategic importance, revenue impact, criticality
2. Adoption Rate -- User adoption %, utilization signals
3. Integration Depth -- API dependencies, systemic coupling
4. Vendor Support -- Roadmap clarity, vendor health, EOL status
5. Cost Efficiency -- Cost-per-user vs market benchmarks
6. Technical Health -- Modernity, tech debt, maintenance burden
7. Risk Score -- Security, compliance, obsolescence, vendor lock-in (higher = more risky)

TELECOM FILTERS (apply when telecom/TMT context): Latency sensitivity, 24x7 SLA, TRAI/GDPR compliance, transaction volume at scale, customer SLA impact.

BEHAVIOR: Be data-driven. State assumptions explicitly. Tag recommendations with confidence: High/Medium/Low. Provide trade-offs, not just conclusions. Use executive-grade language. Never give generic recommendations.

OUTPUT: Always produce structured recommendations with rationale, impact analysis, roadmap (0-3 months | 3-12 months | 12-24 months), and expected outcomes."""


# ═══════════════════════════════════════════════════════════════════════════
# SCORING ENGINE
# ═══════════════════════════════════════════════════════════════════════════
CAT_VALUE = {
    "BSS":9,"OSS":9,"CRM":8.5,"ERP":8.5,"Security":8.5,"Network":8,"Database":8,
    "Cloud":7.5,"Analytics":7.5,"ITSM":7,"Monitoring":6.5,"Logging":6.5,"APM":6.5,
    "DevOps":6,"Collaboration":5.5,"Storage":6,"Other":4.5,
}
CRIT_MOD  = {"Critical":2,"High":1,"Medium":0,"Low":-1.5,None:0}
TIER1_VND = {"microsoft","aws","amazon","google","oracle","sap","salesforce","servicenow",
             "splunk","dynatrace","datadog","ibm","cisco","crowdstrike","elastic","palo alto"}

def score_tool(t: Dict) -> Dict[str, float]:
    eol  = bool(t.get("end_of_life"))
    vnd  = (t.get("vendor") or "").lower()
    age  = t.get("age_years") or 0
    u    = t.get("user_count") or 0
    ints = t.get("integrations") or 0
    cost = t.get("annual_cost") or 0

    bv = min(10, max(0, CAT_VALUE.get(t.get("category","Other"),4.5) + CRIT_MOD.get(t.get("criticality"),0)))
    ar = 5 if not u else (9.5 if u>=2000 else 9 if u>=1000 else 8 if u>=500 else 7 if u>=200 else
         6 if u>=100 else 5 if u>=50 else 4 if u>=20 else 3 if u>=10 else 2)
    id_= 5 if not ints else (9.5 if ints>=25 else 8.5 if ints>=15 else 7.5 if ints>=10 else
         6.5 if ints>=5 else 5.5 if ints>=3 else 4)
    vs = 1 if eol else (8.5 if any(v in vnd for v in TIER1_VND) else 6.5 if vnd else 5)
    cpu = cost / max(u,1) if cost else 0
    ce  = 5 if not cost else (9.5 if cpu<50 else 8.5 if cpu<200 else 7.5 if cpu<500 else
          6 if cpu<1000 else 4.5 if cpu<3000 else 3 if cpu<8000 else 1.5 if cpu<20000 else 0.5)
    th  = 1 if eol else (6 if not age else 9.5 if age<=1 else 8.5 if age<=3 else 7 if age<=5 else
          5.5 if age<=8 else 3.5 if age<=12 else 2.5 if age<=15 else 1.5)
    risk = 4.0
    if eol: risk += 3.5
    if t.get("compliance_required"): risk += 1
    if age > 12: risk += 2
    elif age > 8: risk += 1
    if t.get("criticality") == "Critical": risk += 0.5
    if ints >= 20: risk += 1
    elif ints >= 10: risk += 0.5
    if not vnd: risk += 0.5
    return {"business_value":round(bv,2),"adoption_rate":round(ar,2),"integration_depth":round(id_,2),
            "vendor_support":round(vs,2),"cost_efficiency":round(ce,2),"technical_health":round(th,2),
            "risk_score":round(min(10,max(0,risk)),2)}

def composite(scores: Dict) -> float:
    w = {"business_value":0.25,"adoption_rate":0.15,"integration_depth":0.15,
         "vendor_support":0.15,"cost_efficiency":0.15,"technical_health":0.15}
    val = sum(scores.get(k,5)*wt for k,wt in w.items())
    return round(min(10, max(0, val - max(0,(scores.get("risk_score",5)-5)*0.15))), 2)

def action_6r(scores: Dict, t: Dict) -> str:
    c = composite(scores); risk = scores.get("risk_score",5)
    dep = (t.get("deployment") or "").lower()
    if bool(t.get("end_of_life")) or c < 2.5: return "Retire"
    if c < 3.5 and scores.get("adoption_rate",5) < 3: return "Retire"
    if c < 4.5 and scores.get("cost_efficiency",5) < 3: return "Replace"
    if risk >= 8.5 and c < 6: return "Replace"
    if c >= 7.5 and risk <= 4: return "Retain"
    if c >= 6.5 and risk <= 5: return "Retain"
    if c >= 6 and "on-prem" in dep and scores.get("technical_health",5) >= 5: return "Rehost"
    if c >= 5.5 and scores.get("technical_health",5) < 5.5: return "Replatform"
    if c >= 4 and scores.get("technical_health",5) < 4.5: return "Refactor"
    return "Retain" if c >= 6 else "Replatform"

def confidence(t: Dict) -> str:
    k = sum(1 for f in ["annual_cost","user_count","criticality","vendor","integrations","age_years","deployment"] if t.get(f))
    return "High" if k >= 5 else "Medium" if k >= 3 else "Low"

def apply_scores(tools: List[Dict]) -> List[Dict]:
    for t in tools:
        sc = score_tool(t)
        t.update(scores=sc, composite_score=composite(sc),
                 rationalization_action=action_6r(sc,t), confidence_level=confidence(t))
    return tools


# ═══════════════════════════════════════════════════════════════════════════
# DUPLICATION DETECTOR
# ═══════════════════════════════════════════════════════════════════════════
def detect_dups(tools: List[Dict]) -> List[Dict]:
    by_cat: Dict[str,List] = defaultdict(list)
    for t in tools: by_cat[t.get("category","Other")].append(t)
    dups = []
    for cat, grp in by_cat.items():
        for i in range(len(grp)):
            for j in range(i+1, len(grp)):
                a,b = grp[i],grp[j]
                ov = _ov_score(a,b)
                if ov >= 0.45: dups.append(_dup_rec(a,b,ov,cat))
    return sorted(dups, key=lambda x: x["overlap_percentage"], reverse=True)

def _ov_score(a,b):
    s = 0.35
    if (a.get("subcategory") or "") == (b.get("subcategory") or "") and a.get("subcategory"): s += 0.25
    if (a.get("business_unit") or "").lower() == (b.get("business_unit") or "").lower() and a.get("business_unit"): s += 0.15
    va,vb = (a.get("vendor") or "").lower(),(b.get("vendor") or "").lower()
    if va and vb and va != vb: s += 0.10
    ua,ub = a.get("user_count") or 0, b.get("user_count") or 0
    if ua and ub and min(ua,ub)/max(ua,ub) > 0.5: s += 0.10
    return min(1.0, s)

def _dup_rec(a,b,ov,cat):
    pct = round(ov*100)
    sa,sb = a.get("composite_score",5), b.get("composite_score",5)
    ret,con = (a,b) if sa>=sb else (b,a)
    ca,cb = a.get("annual_cost",0) or 0, b.get("annual_cost",0) or 0
    savings = round(min(ca,cb)*(ov/2))
    return {"id":f"dup-{a['id'][:8]}-{b['id'][:8]}","category":cat,
            "tool_a":a.get("name","?"),"tool_b":b.get("name","?"),"overlap_percentage":pct,
            "retain_candidate":ret.get("name"),"consolidate_candidate":con.get("name"),
            "potential_annual_savings":savings,
            "priority":"High" if pct>=70 else "Medium" if pct>=55 else "Low",
            "rationale":(f"{a.get('name')} and {b.get('name')} show {pct}% functional overlap in {cat}. "
                         f"{ret.get('name')} is the strategic retention candidate "
                         f"(score {max(sa,sb):.1f} vs {min(sa,sb):.1f}). "
                         f"Consolidation could yield ~${savings:,}/yr savings.")}


# ═══════════════════════════════════════════════════════════════════════════
# DATA INGESTION (in-memory)
# ═══════════════════════════════════════════════════════════════════════════
CAT_KW = {
    "Monitoring":["monitor","prometheus","grafana","nagios","zabbix","pagerduty"],
    "Logging":["log","elk","elasticsearch","logstash","kibana","fluentd","splunk","sumo logic"],
    "APM":["apm","trace","dynatrace","newrelic","appdynamics","datadog","jaeger"],
    "Security":["security","siem","firewall","iam","pam","crowdstrike","qualys","veracode","snyk"],
    "ITSM":["itsm","incident","ticket","servicenow","jira service","remedy"],
    "Collaboration":["collab","slack","teams","confluence","sharepoint","zoom","webex"],
    "CRM":["crm","salesforce","hubspot","siebel","dynamics crm"],
    "ERP":["erp","sap","workday","peoplesoft"],
    "BSS":["bss","billing","rating","charging"],
    "OSS":["oss","provisioning","network inventory"],
    "Network":["network","sdwan","cisco","juniper","dns","load balanc"],
    "Cloud":["cloud","aws","azure","gcp","terraform","kubernetes","k8s","openstack"],
    "Analytics":["analytic","bi","tableau","powerbi","qlik","looker","data warehouse"],
    "DevOps":["devops","ci/cd","jenkins","gitlab","github actions","sonar","pipeline"],
    "Storage":["storage","backup","netapp","commvault","s3"],
    "Database":["database","postgres","mysql","mongodb","redis","oracle db","sql server"],
}
COL_ALIAS = {
    "tool_name":"name","application":"name","app_name":"name","platform":"name","tool":"name",
    "vendor_name":"vendor","supplier":"vendor",
    "tool_category":"category","type":"category",
    "annual_license_cost":"annual_cost","cost_usd":"annual_cost","annual_spend":"annual_cost","cost":"annual_cost",
    "users":"user_count","num_users":"user_count","license_count":"user_count","active_users":"user_count",
    "team":"business_unit","department":"business_unit",
    "tool_owner":"owner","contact":"owner",
    "environment":"deployment","hosting":"deployment",
    "integration_count":"integrations","age":"age_years",
    "eol":"end_of_life","deprecated":"end_of_life",
    # additional name aliases
    "application_name":"name","product_name":"name","product":"name",
    "system_name":"name","system":"name","software":"name","solution":"name",
    "service_name":"name","service":"name","asset_name":"name","asset":"name",
    "technology":"name","tech_name":"name","component":"name",
    # additional cost aliases
    "total_cost":"annual_cost","yearly_cost":"annual_cost","license_cost":"annual_cost",
    "license_fee":"annual_cost","subscription_cost":"annual_cost","subscription":"annual_cost",
    "spend":"annual_cost","budget":"annual_cost","annual_budget":"annual_cost",
    "contract_value":"annual_cost","contract_cost":"annual_cost",
    # additional explicit cost aliases
    "software_cost":"annual_cost","tool_cost":"annual_cost","app_cost":"annual_cost",
    "application_cost":"annual_cost","platform_cost":"annual_cost","cost_per_year":"annual_cost",
    "yearly_spend":"annual_cost","annual_expenditure":"annual_cost","expenditure":"annual_cost",
    "annual_fee":"annual_cost","annual_price":"annual_cost","price_per_year":"annual_cost",
    "yearly_fee":"annual_cost","yearly_license":"annual_cost","annual_licence":"annual_cost",
    "licence_cost":"annual_cost","licence_fee":"annual_cost","annual_licence_cost":"annual_cost",
    "saas_cost":"annual_cost","saas_spend":"annual_cost","software_spend":"annual_cost",
    "it_cost":"annual_cost","technology_cost":"annual_cost","tech_cost":"annual_cost",
    "cost_usd_annual":"annual_cost","usd_annual":"annual_cost","annual_value":"annual_cost",
    "total_annual_spend":"annual_cost","total_spend":"annual_cost","total_annual_fee":"annual_cost",
    # additional user aliases
    "seats":"user_count","licenses":"user_count","license_count2":"user_count","headcount":"user_count",
    # additional category aliases
    "tool_type":"category","app_type":"category","type2":"category",
    # additional vendor aliases
    "manufacturer":"vendor","provider":"vendor","publisher":"vendor","oem":"vendor",
    # additional age aliases
    "years_in_use":"age_years","tool_age":"age_years","implementation_year":"age_years",
    # additional integration aliases
    "number_of_integrations":"integrations","num_integrations":"integrations","api_count":"integrations",
    # pain point aliases
    "pain_points":"pain_points","pain_point":"pain_points","challenges":"pain_points",
    "issues":"pain_points","problems":"pain_points","concerns":"pain_points","notes":"pain_points",
    "observations":"pain_points","remarks":"pain_points",
    # org-level IT budget aliases (separate from per-tool annual_cost)
    "it_budget":"it_budget","technology_budget":"it_budget","tech_budget":"it_budget",
    "annual_it_budget":"it_budget","total_it_budget":"it_budget","it_spend":"it_budget",
    "annual_technology_budget":"it_budget","total_technology_budget":"it_budget",
}

def _norm_cat(raw:str)->str:
    low=raw.lower()
    for c,kws in CAT_KW.items():
        if any(k in low for k in kws): return c
    return raw.strip() or "Other"

def _norm_dep(raw:str)->Optional[str]:
    low=raw.lower()
    if any(x in low for x in ["cloud","saas","paas","iaas"]): return "Cloud"
    if any(x in low for x in ["on-prem","onprem","premise"]): return "On-Prem"
    if "hybrid" in low: return "Hybrid"
    return raw or None

def _norm_crit(v)->Optional[str]:
    if not v: return None
    x=str(v).lower()
    if "critical" in x: return "Critical"
    if "high" in x: return "High"
    if "med" in x: return "Medium"
    if "low" in x: return "Low"
    return str(v)

def _s(v:Any)->str:
    return "" if v is None else ("" if str(v).lower() in ("nan","none","null","n/a") else str(v).strip())

def _n(v:Any)->Optional[float]:
    if v is None: return None
    s = str(v).strip().replace(",","").replace("$","").replace("£","").replace("€","").replace(" ","")
    if not s or s.lower() in ("nan","none","null","n/a","-",""):
        return None
    # Handle K / M / B suffixes  (e.g. "50K", "1.2M", "2.5B")
    m = re.match(r'^(-?[\d.]+)\s*([KkMmBb]?)$', s)
    if m:
        num, suffix = float(m.group(1)), m.group(2).upper()
        return num * {"K":1_000,"M":1_000_000,"B":1_000_000_000}.get(suffix, 1)
    try: return float(s)
    except: return None

def _ni(v:Any)->Optional[int]:
    x=_n(v); return int(x) if x is not None else None

def _b(v:Any)->bool:
    return str(v).lower().strip() in ("true","yes","1","y","eol","deprecated","end of life")

def normalize(d:Dict)->Dict:
    pain_raw = _s(d.get("pain_points"))
    t={"id":str(uuid.uuid4()),"name":_s(d.get("name") or d.get("tool_name") or "Unknown Tool"),
       "vendor":_s(d.get("vendor")) or None,"category":_norm_cat(_s(d.get("category",""))),
       "description":_s(d.get("description")) or None,"owner":_s(d.get("owner")) or None,
       "business_unit":_s(d.get("business_unit")) or None,"annual_cost":_n(d.get("annual_cost")),
       "user_count":_ni(d.get("user_count")),"license_type":_s(d.get("license_type")) or None,
       "deployment":_norm_dep(_s(d.get("deployment",""))),
       "criticality":_norm_crit(d.get("criticality")),
       "integrations":_ni(d.get("integrations")),"age_years":_n(d.get("age_years")),
       "end_of_life":_b(d.get("end_of_life",False)),"compliance_required":_b(d.get("compliance_required",False)),
       "pain_points":pain_raw or None}
    return {k:v for k,v in t.items() if v not in (None,"") or k in ("id","name","category","end_of_life")}

def parse_csv_bytes(b:bytes)->List[Dict]:
    import pandas as pd
    df=pd.read_csv(io.BytesIO(b))
    return _df(df)

def parse_excel_bytes(b:bytes)->List[Dict]:
    import pandas as pd
    df=pd.read_excel(io.BytesIO(b))
    return _df(df)

def parse_json_bytes(b:bytes)->List[Dict]:
    data=json.loads(b)
    if isinstance(data,list): return [normalize(t) for t in data]
    for v in data.values():
        if isinstance(v,list): return [normalize(t) for t in v]
    return []

def parse_pdf_bytes(b:bytes)->str:
    try:
        import pdfplumber
        txt=""
        with pdfplumber.open(io.BytesIO(b)) as pdf:
            for pg in pdf.pages: txt+=pg.extract_text() or ""
        return txt[:8000]
    except: return ""

def parse_pptx_bytes(b:bytes)->str:
    try:
        from pptx import Presentation
        prs=Presentation(io.BytesIO(b))
        lines=[]
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        ln=" ".join(run.text for run in para.runs if run.text.strip())
                        if ln.strip(): lines.append(ln.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_txt=" | ".join(c.text.strip() for c in row.cells if c.text.strip())
                        if row_txt: lines.append(row_txt)
        return "\n".join(lines)[:8000]
    except: return ""

async def parse_image_bytes(b:bytes, media_type:str)->List[Dict]:
    import base64
    img_b64=base64.standard_b64encode(b).decode("utf-8")
    resp=await client.messages.create(
        model=ANTHROPIC_MODEL, max_tokens=3000, temperature=0.1,
        messages=[{"role":"user","content":[
            {"type":"image","source":{"type":"base64","media_type":media_type,"data":img_b64}},
            {"type":"text","text":(
                "This image shows a software/tool inventory or IT portfolio. "
                "Extract every tool or application visible. For each tool return: "
                "name, vendor, category, annual_cost (number only, omit currency symbols), "
                "user_count (number only), criticality (Critical/High/Medium/Low), "
                "deployment (Cloud/On-Prem/Hybrid), age_years (number), "
                "integrations_count (number), end_of_life (true/false). "
                'Return ONLY a JSON array: [{"name":"...","vendor":"...",...}, ...]. '
                "Omit fields not visible. Always include name."
            )}
        ]}]
    )
    raw=resp.content[0].text
    try:
        m=re.search(r'\[.*\]', raw, re.DOTALL)
        if m: return [normalize(t) for t in json.loads(m.group())]
    except: pass
    return []

_LAST_INGEST_META: Dict = {}

def _df(df)->List[Dict]:
    import pandas as pd
    df.columns=[c.lower().strip().replace(" ","_").replace("-","_") for c in df.columns]
    # Extract org-level IT budget before renaming (it_budget col is separate from annual_cost)
    detected_budget = None
    for col in df.columns:
        if COL_ALIAS.get(col) == "it_budget":
            vals = df[col].dropna()
            if not vals.empty:
                detected_budget = _n(vals.iloc[0])
            break
    _LAST_INGEST_META["detected_budget"] = detected_budget
    df=df.rename(columns={k:v for k,v in COL_ALIAS.items() if k in df.columns})

    # Fuzzy fallback: detect annual_cost column if still not mapped
    if "annual_cost" not in df.columns:
        COST_KW = ("cost","spend","price","expenditure","licence","license","fee","value","contract","subscription","budget")
        candidates = [c for c in df.columns if any(kw in c for kw in COST_KW)
                      and c not in ("it_budget","technology_budget","tech_budget")]
        for cand in candidates:
            # Accept column if it has at least one parseable numeric value
            clean = df[cand].astype(str).str.replace(r'[\$£€,\s]','',regex=True)
            numeric_count = pd.to_numeric(clean, errors='coerce').notna().sum()
            if numeric_count > 0:
                df = df.rename(columns={cand: "annual_cost"})
                break

    # Fallback: detect name column if not found
    if "name" not in df.columns:
        str_cols = [c for c in df.columns if df[c].dtype == object and c not in
                    ("vendor","category","description","owner","business_unit","deployment",
                     "criticality","license_type","end_of_life","compliance_required","pain_points")]
        if str_cols:
            best = max(str_cols, key=lambda c: df[c].nunique())
            df = df.rename(columns={best: "name"})
    df=df.where(pd.notna(df),None)
    return [normalize(r.to_dict()) for _,r in df.iterrows()]


# ═══════════════════════════════════════════════════════════════════════════
# AI HELPERS
# ═══════════════════════════════════════════════════════════════════════════
def _extract_json(text: str):
    t = text.strip()
    # Strip markdown code fences (```json ... ``` or ``` ... ```)
    t = re.sub(r'^```[a-zA-Z]*\s*\n?', '', t)
    t = re.sub(r'\n?```\s*$', '', t)
    t = t.strip()
    # Find first { or [ in case there's leading text
    for i, ch in enumerate(t):
        if ch in '{[':
            t = t[i:]
            break
    # First attempt: direct parse
    try:
        return json.loads(t)
    except json.JSONDecodeError:
        pass
    # Second attempt: fix literal newlines inside JSON string values
    def _fix_nl(m):
        return m.group(0).replace('\n', '\\n').replace('\r', '')
    try:
        return json.loads(re.sub(r'"(?:[^"\\]|\\.)*"', _fix_nl, t, flags=re.DOTALL))
    except Exception:
        raise ValueError("JSON parse failed")

def _parse_structured_text(text: str) -> List[Dict]:
    """Parse key:value block format (Tool:/Vendor:/Category:/...) without AI."""
    KEY_MAP = {
        "tool": "name", "tool name": "name", "application": "name", "app": "name",
        "vendor": "vendor", "provider": "vendor", "manufacturer": "vendor",
        "category": "category", "type": "category",
        "annual cost": "annual_cost", "cost": "annual_cost", "annual spend": "annual_cost",
        "spend": "annual_cost", "price": "annual_cost", "licence cost": "annual_cost",
        "license cost": "annual_cost",
        "users": "user_count", "user count": "user_count", "seats": "user_count",
        "criticality": "criticality", "priority": "criticality",
        "deployment": "deployment", "hosting": "deployment",
        "integrations": "integrations", "integration count": "integrations",
        "age years": "age_years", "age": "age_years", "years": "age_years",
        "end of life": "end_of_life", "eol": "end_of_life",
        "description": "description", "notes": "description",
        "business unit": "business_unit", "department": "business_unit",
        "compliance": "compliance_required", "compliance required": "compliance_required",
    }
    tools = []
    current: Dict = {}
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            if current.get("name"):
                tools.append(current)
                current = {}
            continue
        if ":" not in line:
            continue
        key_raw, _, val = line.partition(":")
        key = key_raw.strip().lower()
        val = val.strip()
        if not val:
            continue
        mapped = KEY_MAP.get(key)
        if not mapped:
            continue
        if mapped in ("annual_cost", "integrations", "age_years", "user_count"):
            try:
                current[mapped] = int(float(re.sub(r"[^\d.]", "", val)))
            except (ValueError, TypeError):
                pass
        elif mapped == "end_of_life":
            current[mapped] = val.lower() in ("true", "yes", "1", "y")
        elif mapped == "compliance_required":
            current[mapped] = val.lower() in ("true", "yes", "1", "y")
        else:
            current[mapped] = val
    if current.get("name"):
        tools.append(current)
    return tools

async def ai_parse_text(text:str)->List[Dict]:
    resp=await client.messages.create(
        model=ANTHROPIC_MODEL,temperature=0.1,max_tokens=3000,
        messages=[{"role":"user","content":
            f"""Extract all tools/apps/platforms from the text. Return ONLY valid JSON with a "tools" array -- no other text or markdown.
Each item: name, vendor(or null), category(Monitoring/Logging/APM/Security/ITSM/Collaboration/
CRM/ERP/BSS/OSS/Cloud/Analytics/DevOps/Network/Storage/Database/Other),
description(brief/null), annual_cost(number/null), user_count(number/null),
criticality(Critical/High/Medium/Low/null), deployment(Cloud/On-Prem/Hybrid/null),
integrations(number/null), age_years(number/null), end_of_life(bool), compliance_required(bool),
business_unit(string/null).
TEXT:\n{text}"""}])
    try:
        data=_extract_json(resp.content[0].text)
        if isinstance(data,list): return data
        for v in data.values():
            if isinstance(v,list): return v
        return []
    except: return []

async def ai_assess(tools:List[Dict],dups:List[Dict],industry:str,focus:str)->Dict:
    # Collect pain points from tool data
    all_pain = [t["pain_points"] for t in tools if t.get("pain_points")]
    pain_context = f"\nKNOWN PAIN POINTS FROM DATA:\n" + "\n".join(f"- {p}" for p in all_pain[:20]) if all_pain else ""
    resp=await client.messages.create(
        model=ANTHROPIC_MODEL,temperature=0.2,max_tokens=8000,
        system=SYSTEM_PROMPT,
        messages=[{"role":"user","content":f"""Perform a comprehensive rationalization assessment.
Industry: {industry.upper()}{(' | Focus: '+focus) if focus else ''}{pain_context}

TOOL INVENTORY (pre-scored, first 30):
{json.dumps(tools[:30],indent=2)[:8000]}

DUPLICATION ANALYSIS:
{json.dumps(dups[:10],indent=2)[:2000]}

Return ONLY valid JSON (no markdown, no explanation) with these exact keys:
{{
  "executive_summary": {{
    "tool_ecosystem": "Bullet-point list of 4-6 key findings about the current tool ecosystem. Each bullet starts with '- '. Cover: total tool count by category, on-prem vs cloud split, vendor concentration, critical challenges, EOL/legacy exposure. Use specific numbers. Example format: '- 28 tools assessed across 8 categories (APM, Security, ITSM...)\\n- 43% on-premises creating migration debt...'",
    "overlapping_tools": "Bullet-point list of each duplicate/overlap pair. Each bullet starts with '- '. Format each bullet as: '- [Tool A] vs [Tool B] ([Category]) — combined cost $X/yr → Consolidate on [Tool] saving $Y/yr'. End with a total savings bullet. Example: '- Splunk vs ELK (Logging) — combined $225K/yr → Consolidate on ELK saving $180K/yr'",
    "benchmarking": "Bullet-point list of 4-6 benchmarking observations. Each bullet starts with '- '. Compare each flagged tool against market rate or maturity standard. Format: '- [Tool]: $X/user vs market $Y/user — [overpriced/underperforming/EOL]. Recommendation: [action].'",
    "kpis_success_factors": "Bullet-point list of 5-7 KPI comparisons. Each bullet starts with '- '. Use 'Current X → Target Y' format inline. Cover: tool count, annual spend, duplicate pairs, EOL tools, risk level, cloud adoption %, avg portfolio score. Example: '- Tool count: Current 28 → Target 18 (35% reduction)'",
    "recommendations": "Bullet-point list of 6-8 specific platform decisions. Each bullet starts with '- [Action]:'. Group by action type. Format: '- Retain: [Tool] — [rationale]', '- Replace: [Tool A] with [Tool B] — [cost/risk rationale]', '- Retire: [Tool] — [reason]. Estimated savings: $X/yr'",
    "rollout_roadmap": "Structured phase-wise bullet list with dependencies. Format:\\nPhase 1 — Quick Wins (0-3 Months):\\n- [action with dependency note]\\n- [action]\\nPhase 2 — Strategic Delivery (3-12 Months):\\n- [action] (Depends on: Phase 1 [item])\\n- [action]\\nPhase 3 — Transformation (12-24 Months):\\n- [action] (Depends on: Phase 2 [item])\\n- [action]"
  }},
  "three_pillars": {{
    "as_is_assessment": "Bullet-point analysis of 7-9 key findings about the CURRENT AS-IS state of this specific tool portfolio. Each bullet MUST start with '- '. Use the actual tools, costs, categories and scores from the inventory above. Cover: total tool count vs industry benchmark (May 2026 average is 18-25 tools for mid-market), on-prem vs cloud split and migration debt, vendor concentration risk, tools with composite score below 5, EOL/legacy exposure, estimated licence waste from duplicates, integration gaps creating manual work. Cite specific tool names and numbers. Format: '- [Finding with specific tool names or numbers]: [Impact and benchmark context]'",
    "governance_framework": "Bullet-point list of 7-9 governance recommendations tailored to THIS portfolio. Each bullet MUST start with '- '. Apply current (May 2026) best practices: ITIL 4 Service Value Chain, COBIT 2019 governance objectives, ISO/IEC 27001:2022 controls, DORA regulation (if financial sector), and FinOps Foundation cloud spend principles. Cover: tool procurement and approval policy, software licence lifecycle management, shadow IT discovery controls, vendor SLA and contract review cadence, cloud spend guardrails, zero-trust access policy for SaaS tools, rationalization review board cadence, and compliance audit trail requirements. Format: '- [Governance control or policy]: [Specific recommendation citing current framework standard]'",
    "tools_standardization": "Bullet-point list of 7-9 standardization decisions for THIS portfolio. Each bullet MUST start with '- '. Reference current (May 2026) market leaders per category: e.g. security (CrowdStrike Falcon, Wiz, Palo Alto XSIAM), ITSM (ServiceNow, Jira Service Management), observability (Datadog, Grafana, Dynatrace), collaboration (Microsoft 365, Slack), DevOps (GitHub Actions, GitLab CI, Azure DevOps), identity (Entra ID, Okta), data (Snowflake, Databricks, dbt). Cover: category-by-category preferred platform decisions from THIS portfolio, SaaS-first vs on-prem policy, API-first integration standard, single-vendor vs best-of-breed approach per domain, standard evaluation scorecard criteria for new tool requests, sunset timeline for non-standard tools. Format: '- [Category or domain]: [Standard platform decision citing specific tools from the portfolio and market alternatives]'"
  }},
  "portfolio_overview": {{"total_tools":<int>,"total_annual_cost":<float>,"portfolio_health":"Healthy|At Risk|Critical","health_rationale":"<brief>"}},
  "before_after_comparison": {{
    "current_state": {{
      "tool_count": <total number of tools>,
      "annual_cost": <total annual cost as number>,
      "duplicate_pairs": <number of overlapping tool pairs>,
      "eol_tools": <number of end-of-life tools>,
      "risk_level": "High|Medium|Low",
      "key_issues": ["<specific issue naming tools e.g. Splunk+ELK overlap in Logging>","<specific issue>","<specific issue>","<specific issue>"]
    }},
    "future_state": {{
      "projected_tool_count": <estimated count after rationalization>,
      "projected_annual_cost": <estimated annual cost after savings>,
      "estimated_annual_savings": <projected annual savings as number>,
      "eol_tools_resolved": <number of EOL issues to be addressed>,
      "risk_level": "High|Medium|Low",
      "improvements": ["<specific improvement e.g. Retire Symantec, standardise on CrowdStrike>","<specific improvement>","<specific improvement>","<specific improvement>","<specific improvement>"]
    }}
  }},
  "portfolio_pain_areas": ["<specific pain area 1>","<specific pain area 2>","<specific pain area 3>"],
  "rationalization_summary": {{"Retain":<int>,"Rehost":<int>,"Replatform":<int>,"Refactor":<int>,"Replace":<int>,"Retire":<int>}},
  "duplicate_tools": [{{"tool_a":"<name>","tool_b":"<name>","category":"<cat>","overlap_reason":"<why duplicate>","recommendation":"<action>"}}],
  "top_recommendations": [{{"rank":1,"title":"<title>","description":"<detail>","impact":"<impact>","effort":"Low|Medium|High","priority":"Critical|High|Medium","confidence":"High|Medium|Low","timeline":"0-3 months|3-12 months|12-24 months"}}],
  "tool_analysis": [{{"tool_name":"<name>","overview":"<1-2 sentence role in landscape>","facts":[{{"label":"Vendor","value":"<v>"}},{{"label":"Annual Cost","value":"<v>"}},{{"label":"Users","value":"<v>"}},{{"label":"Age","value":"<v>"}},{{"label":"Deployment","value":"<v>"}}],"strengths":["<strength1>","<strength2>"],"challenges":["<challenge1>","<challenge2>"],"gaps":["<gap1>"],"cost_analysis":"<cost vs value narrative>","benchmarking":"<vs market alternatives>","recommendation":"Retain|Rehost|Replatform|Refactor|Replace|Retire"}}],
  "consolidation_opportunities": [{{"tools":["A","B"],"category":"<cat>","overlap_pct":<int>,"recommended_action":"<action>","estimated_savings":<float>,"rationale":"<why>"}}],
  "risk_highlights": [{{"risk_type":"Security|Compliance|Vendor|Obsolescence|Operational","severity":"Critical|High|Medium","affected_tools":["<names>"],"description":"<desc>","mitigation":"<action>"}}],
  "roadmap": {{"short_term":["<0-3m action>"],"medium_term":["<3-12m action>"],"long_term":["<12-24m action>"]}},
  "expected_outcomes": {{"cost_savings_annual":<float>,"risk_reduction":"<desc>","tool_reduction":"<from X to Y>","strategic_gains":"<value>"}}
}}

For tool_analysis, include ALL tools (up to 30). Be specific and data-driven for each section."""}])
    try: return _extract_json(resp.content[0].text)
    except:
        raw = resp.content[0].text
        # Try to extract individual sub-sections
        def _gs(key):
            m = re.search(rf'"{key}"\s*:\s*"((?:[^"\\]|\\.)*)"', raw, re.DOTALL)
            return m.group(1).replace('\\n','\n').replace('\\"','"') if m else ""
        eco   = _gs("tool_ecosystem")
        ovlp  = _gs("overlapping_tools")
        bench = _gs("benchmarking")
        kpis  = _gs("kpis_success_factors")
        rec   = _gs("recommendations")
        road  = _gs("rollout_roadmap")
        p_asis  = _gs("as_is_assessment")
        p_gov   = _gs("governance_framework")
        p_std   = _gs("tools_standardization")
        if eco or ovlp:
            return {"executive_summary": {"tool_ecosystem": eco or "See full assessment.",
                                          "overlapping_tools": ovlp or "",
                                          "benchmarking": bench or "",
                                          "kpis_success_factors": kpis or "",
                                          "recommendations": rec or "",
                                          "rollout_roadmap": road or ""},
                    "three_pillars": {"as_is_assessment": p_asis,
                                      "governance_framework": p_gov,
                                      "tools_standardization": p_std}}
        return {"executive_summary": {"tool_ecosystem": raw[:1000], "overlapping_tools": "", "benchmarking": "", "kpis_success_factors": "", "recommendations": "", "rollout_roadmap": ""}}

def build_report(tools:List[Dict],dups:List[Dict],assessment:Dict,chat_context:Optional[str]=None)->str:
    from datetime import datetime
    total_cost=sum(t.get("annual_cost",0) or 0 for t in tools)
    pot_save=sum(d.get("potential_annual_savings",0) or 0 for d in dups)
    cnts:Dict[str,int]={}
    for t in tools:
        a=t.get("rationalization_action","TBD"); cnts[a]=cnts.get(a,0)+1

    BS={"Retain":"color:#155724;background:#d4edda","Rehost":"color:#004085;background:#cce5ff",
        "Replatform":"color:#856404;background:#fff3cd","Refactor":"color:#7d3c00;background:#fde8d8",
        "Replace":"color:#721c24;background:#f8d7da","Retire":"color:#383d41;background:#e2e3e5",
        "High":"color:#721c24;background:#f8d7da","Medium":"color:#856404;background:#fff3cd",
        "Low":"color:#155724;background:#d4edda","Critical":"color:#fff;background:#c0392b"}
    def bdg(a): return f'<span style="display:inline-block;padding:3px 10px;border-radius:12px;font-size:11px;font-weight:700;{BS.get(a,"color:#333;background:#eee")}">{a}</span>'

    tool_rows="".join(f"<tr><td><b>{t.get('name','')}</b></td><td>{t.get('vendor') or '--'}</td>"
        f"<td>{t.get('category','')}</td><td>{'${:,.0f}'.format(t.get('annual_cost',0)) if t.get('annual_cost') else '--'}</td>"
        f"<td>{t.get('user_count') or '--'}</td><td><b>{t.get('composite_score','--')}</b>/10</td>"
        f"<td>{t.get('scores',{}).get('risk_score','--')}/10</td><td>{bdg(t.get('rationalization_action','TBD'))}</td></tr>" for t in tools)

    dup_rows="".join(f"<tr><td>{d.get('category')}</td><td><b>{d.get('tool_a')}</b></td><td><b>{d.get('tool_b')}</b></td>"
        f"<td><b>{d.get('overlap_percentage')}%</b></td><td>{d.get('retain_candidate')}</td>"
        f"<td>${d.get('potential_annual_savings',0):,.0f}</td><td>{bdg(d.get('priority','Low'))}</td></tr>" for d in dups[:15])

    act_rows="".join(f"<tr><td>{bdg(a)}</td><td><b>{c}</b></td><td>{round(c/max(sum(cnts.values()),1)*100)}%</td></tr>"
        for a,c in cnts.items())

    raw_ex=assessment.get("executive_summary","Run an AI Assessment to generate the executive summary.")
    if isinstance(raw_ex, dict):
        ex_eco   = raw_ex.get("tool_ecosystem","")
        ex_ovlp  = raw_ex.get("overlapping_tools","")
        ex_bench = raw_ex.get("benchmarking","")
        ex_kpis  = raw_ex.get("kpis_success_factors","")
        ex_rec   = raw_ex.get("recommendations","")
        ex_road  = raw_ex.get("rollout_roadmap","")
        ex = None
    else:
        ex = raw_ex
        ex_eco = ex_ovlp = ex_bench = ex_kpis = ex_rec = ex_road = ""
    recs=assessment.get("top_recommendations",[])
    rm=assessment.get("roadmap",{})
    oc=assessment.get("expected_outcomes",{})
    risks=assessment.get("risk_highlights",[])
    pain_areas=assessment.get("portfolio_pain_areas",[])
    dup_tools=assessment.get("duplicate_tools",[])
    tool_analysis=assessment.get("tool_analysis",[])
    bac=assessment.get("before_after_comparison",{})
    pillars=assessment.get("three_pillars",{})

    # Build before/after comparison HTML block
    if bac:
        _cs = bac.get("current_state", {})
        _fs = bac.get("future_state", {})
        def _bac_row(icon, color, text):
            return (f'<div style="display:flex;gap:7px;align-items:flex-start;font-size:12px;'
                    f'padding:5px 0;border-bottom:1px solid rgba(0,0,0,.06)">'
                    f'<span style="color:{color};flex-shrink:0;font-weight:700">{icon}</span>'
                    f'<span style="color:#333">{text}</span></div>')
        _issue_rows = "".join(_bac_row("&#9888;","#E31837",x) for x in (_cs.get("key_issues") or []))
        _impr_rows  = "".join(_bac_row("&#10003;","#00A651",x) for x in (_fs.get("improvements") or []))
        _cs_cost  = "${:,.0f}".format(_cs.get("annual_cost",0) or 0)
        _fs_cost  = "${:,.0f}".format(_fs.get("projected_annual_cost",0) or 0)
        _fs_save  = "${:,.0f}".format(_fs.get("estimated_annual_savings",0) or 0)
        bac_html = f"""<div class="sec pb">
<h2>Before vs After: Rationalization Impact</h2>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:20px">
<div style="background:#fff5f5!important;border:1px solid #fad7d7;border-top:4px solid #E31837;border-radius:8px;padding:20px">
<div style="font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:1px;color:#E31837;margin-bottom:14px">&#128198; CURRENT STATE (AS-IS)</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-bottom:14px">
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:22px;font-weight:800;color:#1a2340">{_cs.get("tool_count","--")}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Tools</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:16px;font-weight:800;color:#1a2340">{_cs_cost}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Annual Cost</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:22px;font-weight:800;color:#E31837">{_cs.get("duplicate_pairs","--")}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Duplicate Pairs</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:16px;font-weight:800;color:#E31837">{_cs.get("risk_level","--")}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Risk Level</div></div>
</div>
<div style="font-size:10px;font-weight:700;color:#E31837;text-transform:uppercase;letter-spacing:.5px;margin-bottom:7px">Key Issues Identified</div>
{_issue_rows or '<div style="font-size:12px;color:#999">No issues listed</div>'}
</div>
<div style="background:#f0faf5!important;border:1px solid #b8e6d0;border-top:4px solid #00A651;border-radius:8px;padding:20px">
<div style="font-size:11px;font-weight:700;text-transform:uppercase;letter-spacing:1px;color:#00A651;margin-bottom:14px">&#9989; FUTURE STATE (POST-RATIONALIZATION)</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-bottom:14px">
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:22px;font-weight:800;color:#00A651">{_fs.get("projected_tool_count","--")}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Tools (Projected)</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:16px;font-weight:800;color:#00A651">{_fs_cost}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Annual Cost (Est.)</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:22px;font-weight:800;color:#00A651">{_fs_save}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Annual Savings</div></div>
<div style="background:#fff!important;padding:10px;border-radius:6px;text-align:center"><div style="font-size:16px;font-weight:800;color:#00A651">{_fs.get("risk_level","--")}</div><div style="font-size:10px;color:#666;text-transform:uppercase;margin-top:2px">Target Risk</div></div>
</div>
<div style="font-size:10px;font-weight:700;color:#00A651;text-transform:uppercase;letter-spacing:.5px;margin-bottom:7px">Improvements After Rationalization</div>
{_impr_rows or '<div style="font-size:12px;color:#999">No improvements listed</div>'}
</div>
</div>
</div>"""
    else:
        bac_html = ""

    # Helper: convert bullet-point text to HTML list
    def _bullets_to_html(text: str) -> str:
        if not text:
            return ""
        lines = text.strip().split("\n")
        html_parts = []
        in_ul = False
        for line in lines:
            stripped = line.strip()
            if stripped.startswith("- "):
                if not in_ul:
                    html_parts.append('<ul style="margin:0;padding-left:18px;list-style:none">')
                    in_ul = True
                item = stripped[2:]
                html_parts.append(
                    f'<li style="padding:5px 0;border-bottom:1px solid rgba(0,0,0,.05);font-size:12px;'
                    f'line-height:1.65;color:#1a2340;display:flex;gap:8px;align-items:flex-start">'
                    f'<span style="color:#0063DC;font-weight:700;flex-shrink:0;margin-top:1px">&#8227;</span>'
                    f'<span>{item}</span></li>'
                )
            else:
                if in_ul:
                    html_parts.append("</ul>")
                    in_ul = False
                if stripped:
                    html_parts.append(f'<p style="font-size:12px;font-weight:700;color:#003366;margin:10px 0 4px">{stripped}</p>')
        if in_ul:
            html_parts.append("</ul>")
        return "".join(html_parts)

    # Helper: render rollout_roadmap with phase sections highlighted
    def _roadmap_panel(text: str) -> str:
        if not text:
            return ""
        PHASE_COLORS = {"Phase 1": "#00A651", "Phase 2": "#0063DC", "Phase 3": "#7B2FBE"}
        PHASE_BG     = {"Phase 1": "#f0faf5",  "Phase 2": "#eef5ff",  "Phase 3": "#f7f0ff"}
        PHASE_ICONS  = {"Phase 1": "1", "Phase 2": "2", "Phase 3": "3"}
        lines = text.strip().split("\n")
        html_parts = []
        current_phase = None
        phase_items = []

        def _flush_phase():
            nonlocal current_phase, phase_items
            if current_phase is None:
                return
            col = PHASE_COLORS.get(current_phase, "#0063DC")
            bg  = PHASE_BG.get(current_phase, "#f8f9ff")
            ico = PHASE_ICONS.get(current_phase, "?")
            items_html = "".join(
                f'<li style="padding:5px 0;border-bottom:1px solid rgba(0,0,0,.05);font-size:12px;'
                f'line-height:1.6;color:#333;display:flex;gap:8px;align-items:flex-start">'
                f'<span style="width:6px;height:6px;border-radius:50%;background:{col};flex-shrink:0;margin-top:5px"></span>'
                f'<span>{item}</span></li>'
                for item in phase_items
            )
            html_parts.append(
                f'<div style="background:{bg}!important;border:1px solid rgba(0,0,0,.08);border-top:3px solid {col};'
                f'border-radius:8px;padding:14px;margin-bottom:10px">'
                f'<div style="display:flex;align-items:center;gap:9px;margin-bottom:9px">'
                f'<div style="width:26px;height:26px;border-radius:50%;background:{col};color:#fff;'
                f'display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:800;flex-shrink:0">{ico}</div>'
                f'<div style="font-size:12px;font-weight:700;color:{col}">{current_phase}</div></div>'
                f'<ul style="margin:0;padding:0;list-style:none">{items_html}</ul></div>'
            )
            phase_items = []
            current_phase = None

        for line in lines:
            stripped = line.strip()
            matched = False
            for ph in PHASE_COLORS:
                if stripped.lower().startswith(ph.lower()):
                    _flush_phase()
                    current_phase = ph
                    matched = True
                    break
            if not matched:
                if stripped.startswith("- "):
                    phase_items.append(stripped[2:])
                elif stripped and current_phase is None:
                    html_parts.append(f'<p style="font-size:12px;color:#555;margin-bottom:6px">{stripped}</p>')
        _flush_phase()
        return "".join(html_parts)

    # Pre-build 6-panel executive summary HTML
    def _ex_panel(bg, border, label_color, label, content, is_roadmap=False):
        if not content:
            return ""
        body = _roadmap_panel(content) if is_roadmap else _bullets_to_html(content)
        if not body:
            body = f'<div style="font-size:12px;line-height:1.8;color:#1a2340;white-space:pre-wrap">{content}</div>'
        return (
            f'<div style="background:{bg}!important;border-left:4px solid {border};border-radius:0 8px 8px 0;padding:14px 18px">'
            f'<div style="font-size:10px;font-weight:700;text-transform:uppercase;letter-spacing:1px;color:{label_color};margin-bottom:10px">{label}</div>'
            f'{body}</div>'
        )
    ex_struct_html = (
        f'<div style="display:flex;flex-direction:column;gap:12px">'
        + _ex_panel("#f0f4fa","#003366","#003366","Current State of Tool Ecosystem &amp; Challenges",ex_eco)
        + _ex_panel("#fff5f5","#E31837","#c0392b","Overlapping Tools — Cost Reduction Opportunities",ex_ovlp)
        + _ex_panel("#f5f0ff","#7B2FBE","#7B2FBE","Benchmarking Analysis vs Market",ex_bench)
        + _ex_panel("#fff8f0","#FFC200","#b08000","Current KPIs vs To-Be Success Factors",ex_kpis)
        + _ex_panel("#f0faf5","#00A651","#00A651","Recommendations — Choosing the Right Platform",ex_rec)
        + _ex_panel("#f8fbff","#0063DC","#0063DC","Implementation &amp; Rollout Plan Roadmap (with Dependencies)",ex_road,is_roadmap=True)
        + '</div>'
    ) if ex is None else f'<div class="exec">{ex}</div>'

    rec_html="".join(f"""<div class="rc">
<div style="display:flex;align-items:center;gap:9px;margin-bottom:8px">
<div style="width:22px;height:22px;background:#0063DC;color:#fff;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;flex-shrink:0">{r.get('rank','')}</div>
<strong style="color:#003366;flex:1;font-size:13px">{r.get('title','')}</strong>
{bdg(r.get('priority','Medium'))} {bdg(r.get('confidence','Medium'))}</div>
<p style="font-size:12px;color:#444;line-height:1.6;margin-bottom:5px">{r.get('description','')}</p>
<p style="font-size:11px;color:#0063DC;font-style:italic">Impact: {r.get('impact','')} &nbsp;&bull;&nbsp; {r.get('timeline','')}</p></div>""" for r in recs[:5])

    def _ph_list(key):
        items=rm.get(key,[])
        return [items] if isinstance(items,str) else items

    _PCFG=[
        ("#00A651","#f0faf5","#b8e6d0","Phase 1","Quick Wins","0 – 3 Months","short_term"),
        ("#0063DC","#eef5ff","#b8d4f8","Phase 2","Strategic Delivery","3 – 12 Months","medium_term"),
        ("#7B2FBE","#f7f0ff","#d8b8f8","Phase 3","Transformation","12 – 24 Months","long_term"),
    ]
    def _phase_col(idx,col,bg,bdr,lbl,title,timeframe,key):
        items=_ph_list(key)
        steps="".join(
            f'<div class="rdmap-item"><div class="rdmap-step" style="background:{col}">{j+1}</div>'
            f'<div class="rdmap-text">{item}</div></div>'
            for j,item in enumerate(items)
        ) or '<div style="font-size:11px;color:#999;padding:8px 0">No items listed</div>'
        return (
            f'<div class="rdmap-col">'
            f'<div class="rdmap-top"><div class="rdmap-node" style="background:{col};box-shadow:0 0 0 5px {col}2a">{idx}</div></div>'
            f'<div class="rdmap-lbl" style="color:{col}">{lbl}</div>'
            f'<div class="rdmap-title">{title}</div>'
            f'<div class="rdmap-time">{timeframe}</div>'
            f'<div class="rdmap-card" style="background:{bg}!important;border-color:{bdr};border-top:3px solid {col}">'
            f'{steps}</div></div>'
        )
    rdmap_cols="".join(_phase_col(i+1,*cfg) for i,cfg in enumerate(_PCFG))
    rdmap_html=(f'<div class="rdmap"><div class="rdmap-track"><div class="rdmap-line"></div>'
                f'<div class="rdmap-phases">{rdmap_cols}</div></div></div>') if rm else ""

    sev_col = {"Critical":"#c0392b","High":"#e67e22","Medium":"#f39c12"}
    risk_html="".join(f"""<div class="rsk">
<div style="font-size:10px;font-weight:700;color:#fff;background:{sev_col.get(r.get('severity','Medium'),'#888')};padding:3px 8px;border-radius:4px;align-self:flex-start;white-space:nowrap">{r.get('severity','').upper()}</div>
<div><div style="font-weight:700;color:#1a2340;font-size:12px">{r.get('risk_type','')} Risk</div>
<div style="font-size:11px;color:#555;margin:3px 0;line-height:1.5">{r.get('description','')}</div>
<div style="font-size:10px;color:#00A651;font-weight:600">Mitigation: {r.get('mitigation','')}</div></div></div>""" for r in risks)

    # ── Reference-style pre-builds ─────────────────────────────────────────
    cnt_r: Dict[str,int] = {}
    for t in tools:
        a=t.get("rationalization_action","TBD"); cnt_r[a]=cnt_r.get(a,0)+1
    retire_eol=sum(1 for t in tools if t.get("end_of_life") or t.get("rationalization_action")=="Retire")
    avg_score=round(sum(t.get("composite_score",0) for t in tools)/max(len(tools),1),1)
    def _sc(s): return "#00A651" if s>=7 else ("#FFC200" if s>=5 else "#E31837")
    ACL={"Retain":"color:#155724;background:#d4edda","Rehost":"color:#004085;background:#cce5ff",
         "Replatform":"color:#856404;background:#fff3cd","Refactor":"color:#7d3c00;background:#fde8d8",
         "Replace":"color:#721c24;background:#f8d7da","Retire":"color:#383d41;background:#e2e3e5","TBD":"color:#333;background:#eee"}
    CMAP={"Retain":"#00A651","Rehost":"#0063DC","Replatform":"#FFC200","Refactor":"#FF7700","Replace":"#E31837","Retire":"#6B7A99","TBD":"#aaa"}
    sav_pct=round(min(100,pot_save/max(total_cost,1)*100))
    kpi_badges="".join(f'<span style="display:inline-block;padding:1px 8px;border-radius:10px;font-size:10px;font-weight:700;{ACL[a]};margin-right:3px;margin-top:5px">{a[:2].upper()}·{cnt_r[a]}</span>' for a in ["Retain","Replace","Retire"] if cnt_r.get(a))
    kpi_html=(
        f'<div class="kpi-row">'
        f'<div class="kc"><div class="kl">TOOLS ASSESSED <span class="kli">i</span></div><div class="kv">{len(tools)}</div><div class="ks">portfolio active</div>{kpi_badges}</div>'
        f'<div class="kc"><div class="kl">ANNUAL SPEND <span class="kli">i</span></div><div class="kv">${total_cost:,.0f}</div><div class="ks">technology budget</div><div class="kbar"><div class="kbf" style="width:70%;background:#003366"></div></div></div>'
        f'<div class="kc"><div class="kl">OVERLAP PAIRS <span class="kli">i</span></div><div class="kv">{len(dups)}</div><div class="ks">consolidation alerts</div></div>'
        f'<div class="kc"><div class="kl">EST. SAVINGS <span class="kli">i</span></div><div class="kv">${pot_save:,.0f}</div><div class="ks">annual potential</div><div class="kbar"><div class="kbf" style="width:{sav_pct}%;background:#00A651"></div></div></div>'
        f'<div class="kc"><div class="kl">RETIRE / EOL <span class="kli">i</span></div><div class="kv" style="color:#E31837">{retire_eol}</div><div class="ks">immediate action</div></div>'
        f'<div class="kc"><div class="kl">AVG TOOL SCORE <span class="kli">i</span></div><div class="kv" style="color:{_sc(avg_score)}">{avg_score}</div><div class="ks">portfolio health</div><div class="kbar"><div class="kbf" style="width:{round(avg_score*10)}%;background:{_sc(avg_score)}"></div></div></div>'
        f'</div>'
    )
    layer_tabs="".join(f'<span class="lt">{a} <span class="lc">{cnt_r.get(a,0)}</span></span>' for a in ["Retain","Rehost","Replatform","Refactor","Replace","Retire"] if cnt_r.get(a,0)>0)
    layer_nav_html=(f'<div class="lnav"><div class="lnav-title">Rationalization Layer Navigator <span style="font-size:10px;font-weight:400;opacity:.5">6R &#8594; Roadmap</span></div>'
                    f'<span class="lt on">All Tools <span class="lc">{len(tools)}</span></span>{layer_tabs}</div>')
    tf_rows="".join(
        f'<div class="tf-item"><div class="tf-r1"><span class="bdg" style="{ACL.get(t.get("rationalization_action","TBD"),ACL["TBD"])}">{t.get("rationalization_action","TBD")}</span>'
        f'<span class="tf-name">{t.get("name","")}</span><span class="tf-score" style="color:{_sc(t.get("composite_score",0))}">{t.get("composite_score",0)}/10</span></div>'
        f'<div class="tf-r2">{t.get("vendor","") or "—"} &bull; {t.get("category","")}'
        f'{(" &bull; ${:,.0f}/yr".format(t.get("annual_cost"))) if t.get("annual_cost") else ""}'
        f'{(" &bull; "+str(t.get("user_count"))+" users") if t.get("user_count") else ""}</div></div>'
        for t in tools[:20])
    tf_html=(f'<div class="tool-feed"><div class="col-hdr"><span class="col-hdr-title">Tool Portfolio</span>'
             f'<span style="font-size:10px;color:#6B7A99">{len(tools)} tools active</span></div>{tf_rows}</div>')
    cost_by_act: Dict[str,float]={}
    for t in tools:
        a2=t.get("rationalization_action","TBD"); cost_by_act[a2]=cost_by_act.get(a2,0.0)+(t.get("annual_cost") or 0)
    _tc2=sum(cost_by_act.values()) or 1.0; _cum2=0.0; ib_rows=""
    for a2,cv in sorted(cost_by_act.items(),key=lambda x:x[1],reverse=True):
        p2=cv/_tc2*100; _cum2+=p2
        ib_rows+=(f'<div class="ib-row"><div class="ib-label">{a2}</div><div class="ib-bar">'
                  f'<div style="height:100%;display:flex;border-radius:4px;overflow:hidden">'
                  f'<div style="width:{p2:.1f}%;background:{CMAP.get(a2,"#aaa")};height:100%"></div>'
                  f'<div style="width:{min(_cum2,100.0)-p2:.1f}%;background:#B8D4EE;height:100%"></div>'
                  f'</div></div><div class="ib-pct">{p2:.0f}% &middot; cum {_cum2:.0f}%</div></div>')
    ib_html=(f'<div class="impact-col"><div class="col-hdr"><span class="col-hdr-title">Impact-Weighted Cost Distribution</span>'
             f'<div style="font-size:10px;color:#666;display:flex;gap:7px;align-items:center">'
             f'<span style="width:8px;height:8px;background:#003366;border-radius:2px;display:inline-block"></span>Value '
             f'<span style="width:8px;height:8px;background:#B8D4EE;border-radius:2px;display:inline-block;margin-left:3px"></span>Cumul.</div></div>{ib_rows}</div>')
    cats_h: Dict[str,Any]={}
    for t in tools:
        cat=t.get("category","Other")
        if cat not in cats_h: cats_h[cat]={"cnt":0,"scores":[],"action":0}
        cats_h[cat]["cnt"]+=1
        if t.get("composite_score"): cats_h[cat]["scores"].append(t["composite_score"])
        if t.get("rationalization_action") in ("Retire","Replace"): cats_h[cat]["action"]+=1
    def _hp(cat:str,d:Dict)->str:
        sc=round(sum(d["scores"])/len(d["scores"]),1) if d["scores"] else 0; col=_sc(sc)
        sev=(f'<div class="hp-sev" style="color:#E31837">Action needed: {d["action"]}</div>') if d["action"] else ""
        return (f'<div class="health-panel"><div class="hp-row">'
                f'<div class="hp-ring" style="background:{col}22;color:{col};border:2px solid {col}">{sc}</div>'
                f'<div style="flex:1"><div class="hp-name">{cat}</div>'
                f'<div class="hp-sub">{d["cnt"]} tool{"s" if d["cnt"]!=1 else ""} &bull; Avg {sc}/10</div>{sev}</div></div></div>')
    hp_html=(f'<div><div class="col-hdr"><span class="col-hdr-title">Category Health</span></div>'
             +"".join(_hp(c,d) for c,d in sorted(cats_h.items(),key=lambda x:x[1]["cnt"],reverse=True)[:7])+'</div>')
    three_col_html=f'<div class="three-col">{tf_html}{ib_html}{hp_html}</div>'

    # ── Three Pillars section ──────────────────────────────────────────────────
    def _pillar_card(icon_svg, accent, bg, title, subtitle, content):
        bullets = _bullets_to_html(content) if content else \
            '<p style="font-size:12px;color:#999">Run assessment to generate content.</p>'
        return (
            f'<div style="background:{bg}!important;border:1px solid rgba(0,0,0,.08);'
            f'border-top:4px solid {accent};border-radius:10px;padding:20px">'
            f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:14px">'
            f'<div style="width:44px;height:44px;background:{accent}18;border-radius:10px;'
            f'display:flex;align-items:center;justify-content:center;flex-shrink:0">'
            f'{icon_svg}</div>'
            f'<div><div style="font-size:14px;font-weight:800;color:#003366;line-height:1.2">{title}</div>'
            f'<div style="font-size:10px;color:#6B7A99;margin-top:2px;text-transform:uppercase;'
            f'letter-spacing:.6px;font-weight:600">{subtitle}</div></div></div>'
            f'{bullets}</div>'
        )

    _icon_asis = ('<svg width="24" height="24" viewBox="0 0 44 44" fill="none" '
                  'xmlns="http://www.w3.org/2000/svg">'
                  '<path d="M22 6 A16 16 0 1 1 8 28" stroke="#1E49E2" stroke-width="2.8" stroke-linecap="round"/>'
                  '<path d="M6 22 L8 28 L14 26" stroke="#1E49E2" stroke-width="2.8" stroke-linecap="round" stroke-linejoin="round"/>'
                  '<path d="M22 38 A16 16 0 0 1 36 16" stroke="#1E49E2" stroke-width="2.8" stroke-linecap="round"/>'
                  '<path d="M38 22 L36 16 L30 18" stroke="#1E49E2" stroke-width="2.8" stroke-linecap="round" stroke-linejoin="round"/>'
                  '<circle cx="22" cy="22" r="5" stroke="#1E49E2" stroke-width="2"/>'
                  '<path d="M22 15v2M22 27v2M15 22h2M27 22h2" stroke="#1E49E2" stroke-width="1.8" stroke-linecap="round"/>'
                  '</svg>')
    _icon_gov  = ('<svg width="24" height="26" viewBox="0 0 40 44" fill="none" '
                  'xmlns="http://www.w3.org/2000/svg">'
                  '<rect x="3" y="2" width="34" height="40" rx="3" stroke="#1E49E2" stroke-width="2.5"/>'
                  '<rect x="8" y="11" width="6" height="6" rx="1" stroke="#1E49E2" stroke-width="2"/>'
                  '<polyline points="9,14 11,16 13,12" stroke="#1E49E2" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"/>'
                  '<line x1="18" y1="14" x2="32" y2="14" stroke="#1E49E2" stroke-width="2" stroke-linecap="round"/>'
                  '<rect x="8" y="20" width="6" height="6" rx="1" stroke="#1E49E2" stroke-width="2"/>'
                  '<polyline points="9,23 11,25 13,21" stroke="#1E49E2" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"/>'
                  '<line x1="18" y1="23" x2="32" y2="23" stroke="#1E49E2" stroke-width="2" stroke-linecap="round"/>'
                  '<rect x="8" y="29" width="6" height="6" rx="1" stroke="#1E49E2" stroke-width="2"/>'
                  '<polyline points="9,32 11,34 13,30" stroke="#1E49E2" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"/>'
                  '<line x1="18" y1="32" x2="32" y2="32" stroke="#1E49E2" stroke-width="2" stroke-linecap="round"/>'
                  '</svg>')
    _icon_std  = ('<svg width="24" height="24" viewBox="0 0 44 44" fill="none" '
                  'xmlns="http://www.w3.org/2000/svg">'
                  '<path d="M18 20 C18 20 20 15 22 12 C23 10 25 10 25 13 L25 18 L31 18 '
                  'C33 18 34 20 33 22 L30 32 C29.5 33.5 28 34 27 34 L16 34 L16 20 Z" '
                  'stroke="#1E49E2" stroke-width="2.2" stroke-linejoin="round"/>'
                  '<rect x="10" y="20" width="6" height="14" rx="1.5" stroke="#1E49E2" stroke-width="2.2"/>'
                  '<circle cx="35" cy="37" r="4" stroke="#1E49E2" stroke-width="1.8"/>'
                  '<path d="M35 33v1M35 40v1M31 37h1M38 37h1" stroke="#1E49E2" stroke-width="1.5" stroke-linecap="round"/>'
                  '<circle cx="39" cy="30" r="3" stroke="#1E49E2" stroke-width="1.6"/>'
                  '<path d="M39 27v1M39 32v1M36 30h1M41 30h1" stroke="#1E49E2" stroke-width="1.4" stroke-linecap="round"/>'
                  '</svg>')

    # ── AI Advisor chat section ───────────────────────────────────────────────
    chat_html = ""
    if chat_context and chat_context.strip():
        def _chat_row(line: str) -> str:
            is_q = line.startswith("Client question:")
            label = line[:line.index(":")+1] if ":" in line else ""
            body  = line[len(label):].strip()
            bg    = "#f0f6ff" if is_q else "#f0faf5"
            border= "#1E49E2" if is_q else "#00A651"
            icon  = "&#10067;" if is_q else "&#129302;"
            return (
                f'<div style="background:{bg}!important;border-left:3px solid {border};'
                f'border-radius:0 6px 6px 0;padding:10px 14px;margin-bottom:8px">'
                f'<div style="font-size:10px;font-weight:700;text-transform:uppercase;'
                f'letter-spacing:.6px;color:{border};margin-bottom:5px">{icon}&nbsp;{label}</div>'
                f'<div style="font-size:12px;line-height:1.65;color:#1a2340">{body}</div></div>'
            )
        rows = "".join(_chat_row(l) for l in chat_context.strip().split("\n") if l.strip())
        if rows:
            chat_html = (
                f'<div class="sec pb" style="page-break-inside:avoid">'
                f'<h2>AI Advisor Session &mdash; Key Insights</h2>'
                f'<p style="font-size:12px;color:#6B7A99;margin-bottom:16px">'
                f'The following questions and answers were generated during the AI Advisor session '
                f'and have been incorporated into this assessment.</p>'
                f'<hr style="border:none;border-top:2px solid #e0e8f4;margin-bottom:16px">'
                f'{rows}</div>'
            )

    pillars_html = ""
    if pillars:
        p_asis = pillars.get("as_is_assessment", "")
        p_gov  = pillars.get("governance_framework", "")
        p_std  = pillars.get("tools_standardization", "")
        if p_asis or p_gov or p_std:
            pillars_html = (
                f'<div class="sec pb">'
                f'<h2>Strategic Pillars: Assessment, Governance &amp; Standardization</h2>'
                f'<p style="font-size:12px;color:#6B7A99;margin-bottom:16px">'
                f'Current-state analysis and forward-looking recommendations across three strategic lenses. '
                f'Benchmarked against May 2026 industry standards and market-leading platforms.</p>'
                f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:16px">'
                + _pillar_card(_icon_asis, "#1E49E2", "#f0f6ff",
                               "As-is Tools Assessment", "Current State &amp; Gap Analysis", p_asis)
                + _pillar_card(_icon_gov,  "#00338D", "#f0f4fa",
                               "Governance Framework", "Policy, Compliance &amp; Controls", p_gov)
                + _pillar_card(_icon_std,  "#00A651", "#f0faf5",
                               "Tools Standardization", "Platform Decisions &amp; Stack Consolidation", p_std)
                + f'</div></div>'
            )

    return f"""<!DOCTYPE html><html lang="en"><head>
<meta charset="UTF-8"><title>Tool Rationalization Report — {datetime.now().strftime('%B %Y')}</title>
<style>
@page{{size:A4;margin:14mm 14mm 18mm 14mm}}
*{{margin:0;padding:0;box-sizing:border-box;-webkit-print-color-adjust:exact!important;print-color-adjust:exact!important;color-adjust:exact!important}}
body{{font-family:'Helvetica Neue',Helvetica,Arial,sans-serif;background:#f0f4fa;color:#1a2340;font-size:13px;line-height:1.5}}
.w{{max-width:1140px;margin:0 auto;background:#fff;box-shadow:0 0 40px rgba(0,0,0,.1)}}
/* Cover / Header */
.hdr{{background:linear-gradient(135deg,#003366 0%,#0063DC 100%);color:#fff;padding:36px 40px 32px}}
.hdr h1{{font-size:22px;font-weight:700;margin-bottom:5px}}.hdr p{{opacity:.75;font-size:12px;margin-top:3px}}
/* KPI row */
.kpi-row{{display:grid;grid-template-columns:repeat(6,1fr);gap:1px;background:#ccd6e8;border-bottom:2px solid #ccd6e8}}
.kc{{background:#fff;padding:16px 14px;text-align:center}}
.kl{{font-size:9px;color:#6B7A99;text-transform:uppercase;letter-spacing:.7px;font-weight:600;display:flex;align-items:center;gap:4px;justify-content:center}}
.kli{{width:13px;height:13px;background:#e8edf5;border-radius:50%;font-size:8px;display:inline-flex;align-items:center;justify-content:center;color:#6B7A99}}
.kv{{font-size:22px;font-weight:800;color:#003366;margin:5px 0 1px}}
.ks{{font-size:10px;color:#8899bb}}
.kbar{{height:4px;background:#e8edf5;border-radius:2px;margin-top:7px;overflow:hidden}}
.kbf{{height:100%;border-radius:2px}}
/* Layer nav */
.lnav{{background:#f5f8ff;border-bottom:1px solid #dde4f5;padding:0 30px;display:flex;align-items:center;gap:8px;overflow-x:auto}}
.lnav-title{{font-size:10px;font-weight:700;color:#6B7A99;text-transform:uppercase;letter-spacing:.5px;margin-right:8px;white-space:nowrap;padding:12px 0}}
.lt{{padding:10px 12px;font-size:11px;font-weight:600;color:#6B7A99;border-bottom:2px solid transparent;white-space:nowrap;display:inline-flex;align-items:center;gap:4px}}
.lt.on{{color:#0063DC;border-bottom-color:#0063DC}}
.lc{{background:#e8edf5;border-radius:8px;font-size:9px;padding:1px 5px;color:#6B7A99}}
/* Three-column */
.three-col{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:0;border-bottom:1px solid #e4eaf5}}
.col-hdr{{display:flex;align-items:center;justify-content:space-between;padding:12px 16px 8px;border-bottom:1px solid #e8edf5}}
.col-hdr-title{{font-size:12px;font-weight:700;color:#003366}}
/* Tool feed */
.tool-feed{{border-right:1px solid #e4eaf5;overflow:hidden}}
.tf-item{{padding:10px 16px;border-bottom:1px solid #f0f3fa;page-break-inside:avoid}}
.tf-r1{{display:flex;align-items:center;gap:7px;margin-bottom:3px}}
.tf-name{{font-size:12px;font-weight:600;color:#1a2340;flex:1}}
.tf-score{{font-size:11px;font-weight:700}}
.tf-r2{{font-size:10px;color:#6B7A99;padding-left:2px}}
/* Impact col */
.impact-col{{border-right:1px solid #e4eaf5;padding:0}}
.ib-row{{display:flex;align-items:center;gap:8px;padding:8px 16px;border-bottom:1px solid #f0f3fa}}
.ib-label{{font-size:10px;font-weight:700;color:#1a2340;width:80px;flex-shrink:0}}
.ib-bar{{flex:1;height:16px;background:#f0f3fa;border-radius:4px;overflow:hidden}}
.ib-pct{{font-size:9px;color:#6B7A99;width:80px;text-align:right;flex-shrink:0}}
/* Category health */
.health-panel{{padding:10px 14px;border-bottom:1px solid #f0f3fa}}
.hp-row{{display:flex;align-items:center;gap:10px}}
.hp-ring{{width:36px;height:36px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:12px;font-weight:800;flex-shrink:0}}
.hp-name{{font-size:12px;font-weight:700;color:#1a2340}}
.hp-sub{{font-size:10px;color:#6B7A99;margin-top:1px}}
.hp-sev{{font-size:10px;font-weight:600;margin-top:2px}}
/* Sections */
.sec{{padding:26px 40px;border-bottom:1px solid #eef1f8}}
.pb{{page-break-before:always}}
h2{{font-size:15px;font-weight:700;color:#003366;margin-bottom:14px;padding-bottom:8px;border-bottom:2px solid #0063DC;page-break-after:avoid}}
/* Tables */
table{{width:100%;border-collapse:collapse;font-size:12px}}
thead{{display:table-header-group}}
thead th{{background:#003366!important;color:#fff!important;padding:9px 12px;text-align:left;font-size:11px;font-weight:600}}
tbody td{{padding:8px 12px;border-bottom:1px solid #f0f3fa;vertical-align:top}}
tbody tr{{page-break-inside:avoid}}
/* Executive summary */
.exec{{background:#f8fbff!important;border-left:4px solid #0063DC;padding:18px 20px;border-radius:0 7px 7px 0;font-size:13px;line-height:1.8;white-space:pre-wrap}}
/* Roadmap visual */
.rdmap{{padding:8px 0 4px}}
.rdmap-track{{position:relative;padding-top:22px}}
.rdmap-line{{position:absolute;top:43px;left:17%;right:17%;height:4px;background:linear-gradient(90deg,#00A651 0%,#0063DC 50%,#7B2FBE 100%);border-radius:2px;z-index:0}}
.rdmap-phases{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:20px;position:relative;z-index:1}}
.rdmap-col{{display:flex;flex-direction:column;align-items:center}}
.rdmap-top{{display:flex;align-items:center;justify-content:center;margin-bottom:10px}}
.rdmap-node{{width:46px;height:46px;border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:20px;font-weight:900;color:#fff;page-break-inside:avoid}}
.rdmap-lbl{{font-size:9px;font-weight:700;text-transform:uppercase;letter-spacing:1.2px;margin-bottom:3px;margin-top:2px}}
.rdmap-title{{font-size:14px;font-weight:800;color:#1a2340;margin-bottom:4px;text-align:center}}
.rdmap-time{{font-size:10px;color:#6B7A99;margin-bottom:12px;background:#eef2f8!important;padding:3px 12px;border-radius:20px;font-weight:700;letter-spacing:.3px}}
.rdmap-card{{border:1px solid;border-radius:8px;padding:14px 16px;width:100%;page-break-inside:avoid}}
.rdmap-item{{display:flex;gap:9px;align-items:flex-start;padding:9px 0;border-bottom:1px solid rgba(0,0,0,.07)}}
.rdmap-item:last-child{{border-bottom:none}}
.rdmap-step{{width:22px;height:22px;border-radius:50%;color:#fff;font-size:10px;font-weight:800;display:flex;align-items:center;justify-content:center;flex-shrink:0;margin-top:1px}}
.rdmap-text{{font-size:12px;color:#333;line-height:1.6}}
/* Rec cards */
.rc{{background:#f8fbff!important;border:1px solid #dde8f8;border-left:4px solid #0063DC;border-radius:7px;padding:14px;margin-bottom:9px;page-break-inside:avoid}}
/* Risk cards */
.rsk{{display:flex;gap:11px;background:#fff5f5!important;border:1px solid #fad7d7;border-radius:8px;padding:12px;margin-bottom:8px;page-break-inside:avoid}}
/* Tool analysis cards */
.tc{{border:1px solid #dde4ef;border-radius:8px;margin-bottom:16px;overflow:hidden;page-break-inside:avoid}}
.tch{{background:#003366!important;color:#fff!important;padding:12px 18px;display:flex;align-items:center;gap:12px}}
.tf{{display:flex;flex-wrap:wrap;border-bottom:1px solid #dde4ef;background:#fff!important}}
.tfi{{padding:9px 13px;border-right:1px solid #dde4ef;min-width:110px}}
.tg{{display:grid;grid-template-columns:1fr 1fr 1fr}}
.tgc{{padding:11px 14px;border-right:1px solid #dde4ef}}
.tgc:last-child{{border-right:none}}
.tga{{display:grid;grid-template-columns:1fr 1fr;border-top:1px solid #dde4ef}}
.tga div{{padding:11px 14px;border-right:1px solid #dde4ef}}
.tga div:last-child{{border-right:none}}
.tag-label{{font-size:9px;font-weight:700;text-transform:uppercase;letter-spacing:.5px;margin-bottom:5px}}
/* Pain / dup cards */
.pain-card{{background:#fff8f0!important;border:1px solid #fde8c8;border-left:4px solid #FFC200;border-radius:5px;padding:10px 14px;font-size:12px;page-break-inside:avoid}}
.dup-card{{background:#fff5f5!important;border:1px solid #fad7d7;border-radius:7px;padding:12px;page-break-inside:avoid}}
/* Footer */
.ftr{{background:#1a2340!important;color:#8899bb;padding:16px 40px;display:flex;justify-content:space-between;font-size:11px}}
/* Grid helpers */
.g2{{display:grid;grid-template-columns:1fr 1fr;gap:10px}}
/* Badge */
.bdg{{display:inline-block;padding:2px 9px;border-radius:10px;font-size:10px;font-weight:700}}
@media screen{{body{{background:#f0f4fa}}.w{{box-shadow:0 0 40px rgba(0,0,0,.1)}}}}
@media print{{body{{background:#fff!important}}.w{{max-width:100%;box-shadow:none}}.pb{{page-break-before:always}}.tc{{page-break-inside:avoid}}.rm{{grid-template-columns:repeat(3,1fr)}}}}
</style></head>
<body><div class="w">

<!-- COVER HEADER -->
<div class="hdr">
<div style="display:flex;align-items:center;gap:16px;margin-bottom:22px">
<div style="width:48px;height:48px;background:rgba(255,255,255,.15);border-radius:10px;display:flex;align-items:center;justify-content:center;font-size:20px;font-weight:900;color:#fff;flex-shrink:0;letter-spacing:-1px">TR</div>
<div style="border-left:1px solid rgba(255,255,255,.25);padding-left:18px">
<div style="font-size:10px;opacity:.6;text-transform:uppercase;letter-spacing:1.4px;margin-bottom:3px">Advisory Services</div>
<div style="font-size:13px;font-weight:600;opacity:.9">Tool Rationalization Agent</div>
</div>
</div>
<h1>Technology Rationalization Assessment Report</h1>
<p style="margin-top:5px">Enterprise Application &amp; Tools Assessment &nbsp;&bull;&nbsp; AI-Powered Advisory &nbsp;&bull;&nbsp; Rationalization Framework</p>
<p style="margin-top:8px;opacity:.55;font-size:11px">Generated: {datetime.now().strftime('%d %B %Y')} &nbsp;&bull;&nbsp; CONFIDENTIAL &nbsp;&bull;&nbsp; For Internal Use Only</p>
</div>

<!-- KPI ROW + LAYER NAV + THREE-COL -->
{kpi_html}
{layer_nav_html}
{three_col_html}

<!-- EXECUTIVE SUMMARY -->
<div class="sec pb">
<h2>Executive Summary</h2>
{ex_struct_html}
</div>

{bac_html}

<!-- THREE STRATEGIC PILLARS -->
{pillars_html}

<!-- AI ADVISOR SESSION -->
{chat_html}

<!-- ACTION SUMMARY -->
<div class="sec">
<h2>Rationalization Action Summary</h2>
<table><thead><tr><th>Action</th><th>Count</th><th>% of Portfolio</th></tr></thead>
<tbody>{act_rows}</tbody></table>
</div>

<!-- RECOMMENDATIONS -->
{'<div class="sec pb"><h2>Top Priority Recommendations</h2>'+rec_html+'</div>' if rec_html else ''}

<!-- ROADMAP -->
{f'<div class="sec pb"><h2>Implementation &amp; Rollout Plan Roadmap</h2>{rdmap_html}</div>' if rdmap_html else ''}

<!-- RISK HIGHLIGHTS -->
{f'<div class="sec"><h2>Risk Highlights</h2>'+risk_html+'</div>' if risk_html else ''}

<!-- PAIN AREAS -->
{(f'<div class="sec"><h2>Portfolio Pain Areas</h2><div class="g2">'+''.join(f'<div class="pain-card">{pa}</div>' for pa in pain_areas)+'</div></div>') if pain_areas else ''}

<!-- DUPLICATE TOOLS -->
{(f'<div class="sec"><h2>Duplicate &amp; Redundant Tools</h2><div class="g2">'+''.join(f'<div class="dup-card"><div style="font-weight:700;color:#003366;margin-bottom:4px;font-size:13px">{d.get("tool_a","")} &harr; {d.get("tool_b","")}</div><div style="font-size:11px;color:#666;margin-bottom:5px"><b>{d.get("category","")}</b> &mdash; {d.get("overlap_reason","")}</div><div style="font-size:11px;font-weight:700;color:#0063DC">&#9654; {d.get("recommendation","")}</div></div>' for d in dup_tools[:10])+'</div></div>') if dup_tools else ''}

<!-- PER-TOOL ANALYSIS -->
{(f'<div class="sec pb"><h2>Per-Tool Analysis &mdash; Consulting View</h2>'+''.join(f'''<div class="tc">
<div class="tch">
<div style="flex:1"><div style="font-size:14px;font-weight:700">{ta.get("tool_name","")}</div><div style="font-size:11px;opacity:.7;margin-top:2px">{ta.get("overview","")}</div></div>
<span class="bdg" style="background:rgba(255,255,255,.18);color:#fff;font-size:11px;padding:3px 12px">{ta.get("recommendation","")}</span>
</div>
<div class="tf">{"".join(f'<div class="tfi"><div class="tag-label" style="color:#888">{f.get("label","")}</div><div style="font-size:12px;font-weight:600;color:#003366">{f.get("value","--")}</div></div>' for f in (ta.get("facts") or [])[:5])}</div>
<div class="tg">
<div class="tgc"><div class="tag-label" style="color:#00A651">&#10003; Strengths</div>{"".join(f'<div style="font-size:11px;color:#333;padding:3px 0;border-bottom:1px solid #f0f0f0">{s}</div>' for s in (ta.get("strengths") or []))}</div>
<div class="tgc"><div class="tag-label" style="color:#E31837">&#9888; Challenges</div>{"".join(f'<div style="font-size:11px;color:#333;padding:3px 0;border-bottom:1px solid #f0f0f0">{c}</div>' for c in (ta.get("challenges") or []))}</div>
<div class="tgc" style="border-right:none"><div class="tag-label" style="color:#b08000">&#9654; Gaps</div>{"".join(f'<div style="font-size:11px;color:#333;padding:3px 0;border-bottom:1px solid #f0f0f0">{g}</div>' for g in (ta.get("gaps") or []))}</div>
</div>
<div class="tga">
<div><div class="tag-label" style="color:#0063DC">Cost Analysis</div><div style="font-size:11px;color:#444;line-height:1.6">{ta.get("cost_analysis","")}</div></div>
<div><div class="tag-label" style="color:#7B2FBE">Benchmarking</div><div style="font-size:11px;color:#444;line-height:1.6">{ta.get("benchmarking","")}</div></div>
</div>
</div>''' for ta in tool_analysis)+'</div>') if tool_analysis else ''}

<!-- FULL PORTFOLIO TABLE -->
<div class="sec pb">
<h2>Full Tool Portfolio Assessment</h2>
<table><thead><tr><th>Tool</th><th>Vendor</th><th>Category</th><th>Annual Cost</th><th>Users</th><th>Score</th><th>Risk</th><th>Action</th></tr></thead>
<tbody>{tool_rows}</tbody></table>
</div>

<!-- DUPLICATION TABLE -->
{f'''<div class="sec"><h2>Duplication &amp; Consolidation Opportunities</h2>
<table><thead><tr><th>Category</th><th>Tool A</th><th>Tool B</th><th>Overlap</th><th>Retain</th><th>Est. Savings</th><th>Priority</th></tr></thead>
<tbody>{dup_rows}</tbody></table></div>''' if dups else ''}

<!-- OUTCOMES -->
{f'''<div class="sec"><h2>Expected Business Outcomes</h2>
<table><tbody>{"".join(f"<tr><td style='width:40%;font-weight:600'>{k.replace('_',' ').title()}</td><td>{v}</td></tr>" for k,v in oc.items())}</tbody></table></div>''' if oc else ''}

<div class="ftr">
<span>Tool Rationalization Agent &nbsp;&bull;&nbsp; Enterprise Technology Strategy Advisory</span>
<span>CONFIDENTIAL &nbsp;&bull;&nbsp; {datetime.now().strftime('%Y')}</span>
</div>
</div></body></html>"""


# ═══════════════════════════════════════════════════════════════════════════
# PYDANTIC MODELS
# ═══════════════════════════════════════════════════════════════════════════
class ChatReq(BaseModel):
    message: str
    history: List[Dict[str, str]] = []
    tools: List[Dict] = []

class AssessReq(BaseModel):
    tools: List[Dict]
    duplications: List[Dict] = []
    industry: str = "telecom"
    focus: Optional[str] = None

class ReportReq(BaseModel):
    tools: List[Dict]
    duplications: List[Dict] = []
    assessment: Dict = {}
    chat_context: Optional[str] = None


# ═══════════════════════════════════════════════════════════════════════════
# ROUTES
# ═══════════════════════════════════════════════════════════════════════════
@app.post("/api/chat")
async def chat(req: ChatReq):
    system_parts = [SYSTEM_PROMPT]
    if req.tools:
        summary = "\n".join(f"- {t.get('name')} | {t.get('category')} | Action:{t.get('rationalization_action','TBD')} | Score:{t.get('composite_score','?')}" for t in req.tools[:30])
        system_parts.append(f"Current tool inventory ({len(req.tools)} tools):\n{summary}")
    msgs = [{"role":m.get("role","user"),"content":m.get("content","")}
            for m in req.history[-20:] if m.get("role") in ("user","assistant")]
    msgs.append({"role":"user","content":req.message})
    resp = await client.messages.create(
        model=ANTHROPIC_MODEL, system="\n\n".join(system_parts), messages=msgs,
        max_tokens=2000, temperature=0.2)
    return {"reply": resp.content[0].text}


@app.post("/api/ingest")
async def ingest(file: Optional[UploadFile]=File(None), text: Optional[str]=Form(None)):
    tools: List[Dict] = []
    if file and file.filename:
        try:
            content = await file.read()
            ext = file.filename.lower().rsplit(".",1)[-1]
            if   ext == "csv":
                tools = parse_csv_bytes(content)
            elif ext in ("xlsx","xls"):
                tools = parse_excel_bytes(content)
            elif ext == "pdf":
                raw = parse_pdf_bytes(content)
                if raw: tools = [normalize(t) for t in await ai_parse_text(raw)]
            elif ext in ("pptx","ppt"):
                raw = parse_pptx_bytes(content)
                if raw: tools = [normalize(t) for t in await ai_parse_text(raw)]
            elif ext in ("jpg","jpeg"):
                tools = await parse_image_bytes(content, "image/jpeg")
            elif ext == "png":
                tools = await parse_image_bytes(content, "image/png")
            else:
                raise HTTPException(400, f"Unsupported file type: .{ext}. Supported: CSV, Excel, PDF, PPT/PPTX, JPG, PNG.")
        except HTTPException: raise
        except Exception as ex:
            raise HTTPException(400, f"File parsing failed: {str(ex)}. Ensure the file has a header row with tool data.")
    elif text:
        try:
            # Try fast direct parser first (handles Tool:/Vendor:/... block format)
            direct = _parse_structured_text(text)
            if direct:
                tools = [normalize(t) for t in direct]
            else:
                # Fall back to AI for unstructured / prose text
                tools = [normalize(t) for t in await ai_parse_text(text)]
        except HTTPException:
            raise
        except Exception as ex:
            raise HTTPException(500, f"AI text parsing failed: {str(ex)}")
    else:
        raise HTTPException(400, "No file or text provided.")
    tools = apply_scores(tools)
    if not tools:
        raise HTTPException(400, "No tools found in the uploaded data. Check that your file has a header row and tool names.")
    dups  = detect_dups(tools)
    total_cost = sum(t.get("annual_cost",0) or 0 for t in tools)
    # Aggregate explicit pain points from tool data
    explicit_pain = [t["pain_points"] for t in tools if t.get("pain_points")]
    # Generate derived pain points from scored results so Phase 4 is always pre-filled
    auto_pain: List[str] = []
    eol_tools = [t["name"] for t in tools if t.get("end_of_life")]
    if eol_tools:
        auto_pain.append(f"End-of-life / deprecated tools requiring urgent retirement: {', '.join(eol_tools[:4])}")
    high_risk = [t["name"] for t in tools if (t.get("scores") or {}).get("risk_score", 0) >= 7]
    if high_risk:
        auto_pain.append(f"High-risk tools with security or compliance exposure: {', '.join(high_risk[:4])}")
    retire_replace = [t["name"] for t in tools if t.get("rationalization_action") in ("Retire","Replace")]
    if retire_replace:
        auto_pain.append(f"{len(retire_replace)} tools recommended for Retire/Replace: {', '.join(retire_replace[:4])}")
    if dups:
        pot = sum(d.get("potential_annual_savings",0) for d in dups)
        auto_pain.append(f"{len(dups)} overlapping tool pairs detected — consolidation opportunity ~${pot:,.0f}/yr savings")
    low_adopt = [t["name"] for t in tools if (t.get("scores") or {}).get("adoption_rate",5) < 3 and (t.get("annual_cost") or 0) > 0]
    if low_adopt:
        auto_pain.append(f"Low-adoption high-cost tools: {', '.join(low_adopt[:3])}")
    # Merge: explicit first, then auto (deduplicate)
    seen: set = set(explicit_pain)
    pain_points = list(explicit_pain)
    for p in auto_pain:
        if p not in seen:
            pain_points.append(p)
            seen.add(p)
    detected_budget = _LAST_INGEST_META.get("detected_budget")
    return {"tools":tools,"duplications":dups,
            "pain_points": pain_points,
            "detected_budget": detected_budget,
            "summary":{"total_tools":len(tools),
                        "total_annual_cost": total_cost,
                        "duplications_found":len(dups),
                        "potential_savings":sum(d.get("potential_annual_savings",0) for d in dups)}}


@app.post("/api/assess")
async def assess(req: AssessReq):
    if not req.tools: raise HTTPException(400,"No tools provided.")
    try:
        result = await ai_assess(req.tools, req.duplications, req.industry, req.focus or "")
    except HTTPException:
        raise
    except Exception as ex:
        raise HTTPException(500, f"AI assessment failed: {str(ex)}")
    return {"assessment": result}


@app.post("/api/report")
async def report(req: ReportReq):
    return {"html": build_report(req.tools, req.duplications, req.assessment, req.chat_context)}


@app.get("/api/health")
async def health():
    return {"status":"healthy","model":ANTHROPIC_MODEL,"key_set":bool(ANTHROPIC_API_KEY)}
from mangum import Mangum
handler = Mangum(app)
